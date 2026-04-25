from __future__ import annotations

from pathlib import Path
from zipfile import ZipFile
import re
import xml.etree.ElementTree as ET

import pandas as pd
from openpyxl import load_workbook


def list_excel_sheets(path: Path) -> list[str]:
    """Retourne la liste des feuilles d'un classeur Excel."""
    workbook = load_workbook(path, read_only=True, data_only=True)
    try:
        return workbook.sheetnames
    finally:
        workbook.close()


def load_excel_sheet(path: Path, sheet_name: str, header: int | None = None) -> pd.DataFrame:
    """Charge une feuille Excel avec les valeurs calculées des formules."""
    return pd.read_excel(path, sheet_name=sheet_name, header=header, engine="openpyxl")


def load_excel_rows(path: Path, sheet_name: str) -> list[list[object]]:
    """Charge les lignes non vides d'une feuille Excel."""
    workbook = load_workbook(path, read_only=True, data_only=True)
    worksheet = workbook[sheet_name]
    try:
        rows = []
        for row in worksheet.iter_rows(values_only=True):
            values = list(row)
            if any(value not in (None, "") for value in values):
                rows.append(values)
        return rows
    finally:
        workbook.close()


def extract_presentation_text(path: Path) -> list[dict[str, object]]:
    """Extrait le texte d'un PowerPoint avec fallback XML léger."""
    try:
        from pptx import Presentation
    except ModuleNotFoundError:
        return _extract_presentation_text_from_xml(path)

    presentation = Presentation(path)
    slides = []
    for index, slide in enumerate(presentation.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                cleaned = " ".join(shape.text.split())
                if cleaned:
                    texts.append(cleaned)
        slides.append({"slide": index, "text": texts})
    return slides


def _extract_presentation_text_from_xml(path: Path) -> list[dict[str, object]]:
    namespace = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    with ZipFile(path) as archive:
        slide_names = sorted(
            [
                name
                for name in archive.namelist()
                if re.match(r"ppt/slides/slide\d+\.xml$", name)
            ],
            key=lambda name: int(re.search(r"slide(\d+)\.xml", name).group(1)),
        )
        slides = []
        for index, name in enumerate(slide_names, start=1):
            xml_root = ET.fromstring(archive.read(name))
            texts = [
                text_node.text.strip()
                for text_node in xml_root.findall(".//a:t", namespace)
                if text_node.text and text_node.text.strip()
            ]
            slides.append({"slide": index, "text": texts})
    return slides
